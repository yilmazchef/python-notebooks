#!/usr/bin/env python


from progressbar import AnimatedMarker, Bar, BouncingBar, Counter, ETA, \
    AdaptiveETA, FileTransferSpeed, FormatLabel, Percentage, \
    ProgressBar, ReverseBar, RotatingMarker, \
    SimpleProgress, Timer, UnknownLength
import os
import sys
import json
import time
import uuid
from folder_structure import notebook_file_paths
from docx_manager import addFooter, addHeader
import collections.abc
from pptx import Presentation
from pptx.util import Cm, Inches, Pt

c = collections
c.abc = collections.abc


def to_py(ipynb_file):
    py_file = ipynb_file.replace(".ipynb", ".py")

    code = json.load(open(ipynb_file))

    for cell in code['cells']:
        if cell['cell_type'] == 'code':
            print('# -------- code --------')
            for line in cell['source']:
                print(line, end='')
            print('\n')
        elif cell['cell_type'] == 'markdown':
            print('# -------- markdown --------')
            for line in cell['source']:
                print("#", line, end='')
            print('\n')


def to_pptx(ipynb_path, pptx_path, header_image_path):

    prs = Presentation(pptx_path)
    blank_slide_layout = prs.slide_layouts[6]

    code = json.load(open(ipynb_path))

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(1)
    pic = slide.shapes.add_picture(header_image_path, left, top)

    for cell in code['cells']:
        if cell['cell_type'] == 'code':
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "Python Code: "

            for line in cell['source']:
                p = tf.add_paragraph()
                p.text = "``` " + line + " ```"
                p.font.size = Pt(12)

            p.text += "\n"

        elif cell['cell_type'] == 'markdown':
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "Info:"

            for line in cell['source']:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(12)

            p.text += "\n"

    prs.save(pptx_path)


def to_md(ipynb_file, md_file):
    cmdlet = f"jupyter nbconvert --to markdown \"{ipynb_file}\""
    os.system(cmdlet)


def to_docx(md_file, docx_file):
    cmdlet = f"pandoc \"{md_file}\" -f markdown -o \"{docx_file}\""
    os.system(cmdlet)

    text = """
    Rouppeplein 16
    1000 Brussel
    Tel. 02 411 29 07
    ondernemingsnr. 0475319893
    RPR BRUSSEL-NEDERLANDSTALIG
    wouter.vandenberge@intecbrussel.be                                
    """
    image = "Temp/IntecHeader.png"
    addHeader(docx_file, image, text)

    image = "Temp/IntecFooter.png"
    addFooter(docx_file, image)


def to_odt(md_file, odt_file):
    print(f"{md_file} to \"{odt_file}\" conversion started...")
    cmdlet = f"pandoc -t odt \"{md_file}\" -o \"{odt_file}\""
    print(cmdlet)
    os.system(cmdlet)
    print(f"\"{md_file}\" to \"{odt_file}\" conversion complete...")


def to_json(md_file, json_file):
    pass


if __name__ == "__main__":

    src_path = sys.argv[1] if len(sys.argv) > 1 else os.getcwd()
    ipynb_file_list = notebook_file_paths(src_path, ".ipynb")

    pBarMax = len(ipynb_file_list)
    widgets = [Percentage(),
               ' ', Bar(),
               ' ', ETA(),
               ' ', AdaptiveETA()]
    pBar = ProgressBar(widgets=widgets, maxval=pBarMax)

    pBar.start()
    pBarCount = 0
    for ipynb_file in ipynb_file_list:

        md_file = ipynb_file.replace(".ipynb", ".md")
        docx_file = ipynb_file.replace(".ipynb", ".docx")
        pptx_file = ipynb_file.replace(".ipynb", ".pptx")

        to_md(ipynb_file, md_file)
        to_docx(md_file, docx_file)
        to_pptx(ipynb_file, pptx_file, os.path.join(os.getcwd(), "Temp", "IntecHeader.png"))

        pBarCount += 1
        pBar.update(pBarCount)

    pBar.finish()
