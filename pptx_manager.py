import collections.abc
from pptx import Presentation
from pptx.util import Cm, Inches, Pt
import json
import os

c = collections
c.abc = collections.abc


def to_pptx(ipynb_path, header_image_path=None):

    pptx_path = ipynb_path.replace(".ipynb", ".pptx")
    prs = Presentation(pptx_path)
    blank_slide_layout = prs.slide_layouts[6]

    code = json.load(open(ipynb_path))

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(1)
    if header_image_path is not None:
        slide.shapes.add_picture(header_image_path, left, top)

    for cell in code['cells']:
        if cell['cell_type'] == 'code':
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "Python Code: "

            p = tf.add_paragraph()
            p.font.size = Pt(12)
            for line in cell['source']:
                p.text += "``` " + line + " ```"

            p.text += "\n"

        elif cell['cell_type'] == 'markdown':
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = width = height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "Info:"

            p = tf.add_paragraph()
            p.font.size = Pt(12)
            for line in cell['source']:
                p.text += line

            p.text += "\n"

    prs.save(pptx_path)


if __name__ == "__main__":

    ipynb_path = os.path.join(
        os.getcwd(), "Notebooks", "English", "02 - Objects and data structures", "01-Numbers.ipynb"
    )

    pptx_path = os.path.join(
        os.getcwd(), "Notebooks", "English", "02 - Objects and data structures", "01-Numbers.pptx"
    )

    header_image_path = os.path.join(
        os.getcwd(), "Temp", "IntecHeader.png"
    )

    to_pptx(ipynb_path, header_image_path)
